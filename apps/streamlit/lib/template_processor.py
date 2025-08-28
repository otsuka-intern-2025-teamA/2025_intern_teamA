"""
テンプレート処理モジュール - PowerPointテンプレートの変数置換
元のテキストの色、サイズ、フォントを保持しながら変数を置換
"""

import os
import shutil
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class TemplateProcessor:
    """PowerPointテンプレート処理クラス"""
    
    def __init__(self, template_path: str):
        """
        テンプレート処理クラスの初期化
        
        Args:
            template_path: テンプレートファイルのパス
        """
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"テンプレートファイルが見つかりません: {template_path}")
    
    def process_template(
        self, 
        variables: dict[str, str], 
        output_path: str,
        preserve_formatting: bool = True
    ) -> str:
        """
        テンプレートを処理して変数を置換
        
        Args:
            variables: 置換する変数の辞書
            output_path: 出力ファイルのパス
            preserve_formatting: フォーマット保持フラグ
            
        Returns:
            出力ファイルのパス
        """
        # テンプレートをコピー
        output_path = Path(output_path)
        shutil.copy2(self.template_path, output_path)
        
        # プレゼンテーションを開く
        prs = Presentation(output_path)
        
        # 各スライドで変数を置換
        total_replacements = 0
        
        # 変数の確認（簡潔に）
        product_vars = {k: v for k, v in variables.items() if "PRODUCTS" in k}
        print(f"🔍 テンプレート処理: {len(variables)}件の変数、{len(product_vars)}件の製品変数")
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_replacements = self._process_slide(
                slide, variables, preserve_formatting
            )
            print(f"🔍 スライド {slide_idx + 1}: {slide_replacements}件の置換")
            total_replacements += slide_replacements
        
        # 保存
        prs.save(output_path)
        
        print(f"テンプレート処理完了: {total_replacements}件の置換を実行")
        return str(output_path)
    
    def _process_slide(self, slide, variables: dict[str, str], preserve_formatting: bool) -> int:
        """スライド内の変数を処理"""
        replacements = 0
        
        # スライド内の各シェイプを処理
        for shape in slide.shapes:
            replacements += self._process_shape(
                shape, variables, preserve_formatting
            )
        
        return replacements
    
    def _process_shape(self, shape, variables: dict[str, str], preserve_formatting: bool) -> int:
        """シェイプ内の変数を処理"""
        replacements = 0
        
        # グループシェイプの場合は再帰処理
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                replacements += self._process_shape(
                    sub_shape, variables, preserve_formatting
                )
            return replacements
        
        # テキストフレームの処理
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            replacements += self._process_text_frame(
                shape.text_frame, variables, preserve_formatting
            )
        
        # テーブルの処理
        if hasattr(shape, "has_table") and shape.has_table:
            replacements += self._process_table(
                shape.table, variables, preserve_formatting
            )
        
        return replacements
    
    def _process_text_frame(self, text_frame, variables: dict[str, str], preserve_formatting: bool) -> int:
        """テキストフレーム内の変数を処理"""
        if not text_frame or not text_frame.text:
            return 0
        
        replacements = 0
        current_text = text_frame.text  # 現在のテキスト（各置換後に更新）
        
        print(f"🔍 _process_text_frame: '{current_text[:50]}...'")
        
        # 各変数をチェック
        for placeholder, value in variables.items():
            # None値のチェック（空文字列は許可）
            if value is None:
                continue
                
            if placeholder in current_text:
                print(f"✅ 変数置換: {placeholder} → {value[:50]}...")
                # フォーマット保持の場合は特別な処理
                if preserve_formatting:
                    replacements_before = replacements
                    replacements += self._replace_with_formatting(
                        text_frame, placeholder, value
                    )
                    print(f"   _replace_with_formatting 戻り値: {replacements - replacements_before}件の置換")
                    # 置換後に現在のテキストを更新
                    current_text = text_frame.text
                else:
                    # 単純な置換
                    current_text = current_text.replace(placeholder, value)
                    text_frame.text = current_text
                    replacements += 1
                    print(f"   単純な置換: +1")
        
        print(f"🔍 _process_text_frame 完了: {replacements}件の置換")
        return replacements
    
    def _replace_with_formatting(self, text_frame, placeholder: str, value: str) -> int:
        """
        フォーマットを保持しながら変数を置換
        元のテキストの色、サイズ、フォントを保持
        """
        if not text_frame.text or placeholder not in text_frame.text:
            return 0
        
        # None値のチェック（空文字列は許可）
        if value is None:
            return 0
        
        replacement_count = 0
        
        # 段落ごとに処理
        for paragraph_idx, paragraph in enumerate(text_frame.paragraphs):
            if placeholder in paragraph.text:
                # PRODUCTS変数の場合のみ特別な処理（複数ラン対応）
                if "PRODUCTS" in placeholder and len(paragraph.runs) > 1:
                    # 全テキストを結合して置換
                    full_text = paragraph.text
                    if placeholder in full_text:
                        # 最初のランに置換後のテキストを設定
                        paragraph.runs[0].text = value
                        
                        # 残りのランを空文字列に設定（削除できないため）
                        for i in range(1, len(paragraph.runs)):
                            paragraph.runs[i].text = ""
                        
                        replacement_count += 1
                else:
                    # 通常の処理（従来の方法）
                    # フォーマット情報を保存
                    original_font = paragraph.font
                    font_info = {
                        "name": original_font.name,
                        "size": original_font.size,
                        "bold": original_font.bold,
                        "italic": original_font.italic,
                        "underline": original_font.underline,
                        "color": original_font.color.rgb if hasattr(original_font.color, 'rgb') and original_font.color.rgb else None,
                    }
                    
                    # テキストを置換
                    paragraph.text = paragraph.text.replace(placeholder, value)
                    
                    # フォーマットを復元
                    if font_info["name"]:
                        paragraph.font.name = font_info["name"]
                    if font_info["size"]:
                        paragraph.font.size = font_info["size"]
                    if font_info["bold"] is not None:
                        paragraph.font.bold = font_info["bold"]
                    if font_info["italic"] is not None:
                        paragraph.font.italic = font_info["italic"]
                    if font_info["underline"] is not None:
                        paragraph.font.underline = font_info["underline"]
                    if font_info["color"] and hasattr(paragraph.font.color, 'rgb'):
                        paragraph.font.color.rgb = font_info["color"]
                    
                    replacement_count += 1
                    
                print(f"    🔍 _replace_with_formatting 完了: {replacement_count}件の置換")
        return replacement_count
    
    def _process_table(self, table, variables: dict[str, str], preserve_formatting: bool) -> int:
        """テーブル内の変数を処理"""
        replacements = 0
        
        # 製品変数の確認
        product_vars = {k: v for k, v in variables.items() if "PRODUCTS" in k}
        if product_vars:
            print(f"🔍 テーブル処理で使用する製品変数: {list(product_vars.keys())}")
        
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.text_frame and cell.text_frame.text:
                    if "PRODUCTS" in cell.text_frame.text:
                        print(f"🔍 テーブルセル[{row_idx+1},{col_idx+1}]で製品変数を発見: {cell.text_frame.text[:100]}...")
                    replacements += self._process_text_frame(
                        cell.text_frame, variables, preserve_formatting
                    )
        
        return replacements
    
    def get_template_info(self) -> dict[str, Any]:
        """テンプレートの情報を取得"""
        try:
            prs = Presentation(self.template_path)
            
            info = {
                "file_path": str(self.template_path),
                "file_size": self.template_path.stat().st_size,
                "slide_count": len(prs.slides),
                "slides": []
            }
            
            # 各スライドの情報を収集
            for i, slide in enumerate(prs.slides):
                slide_info = {
                    "slide_number": i + 1,
                    "shapes_count": len(slide.shapes),
                    "text_placeholders": []
                }
                
                # テキストプレースホルダーを検索
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        if shape.text_frame.text:
                            # 変数プレースホルダーを検索
                            import re
                            placeholders = re.findall(r'\{\{[^}]+\}\}', shape.text_frame.text)
                            if placeholders:
                                slide_info["text_placeholders"].extend(placeholders)
                                # 製品変数の確認
                                product_placeholders = [p for p in placeholders if "PRODUCTS" in p]
                                if product_placeholders:
                                    print(f"🔍 スライド{i+1}で製品変数を発見: {product_placeholders}")
                
                info["slides"].append(slide_info)
            
            return info
            
        except Exception as e:
            return {
                "error": str(e),
                "file_path": str(self.template_path)
            }
    
    def validate_variables(self, variables: dict[str, str]) -> dict[str, Any]:
        """変数の妥当性を検証"""
        validation_result = {
            "valid": True,
            "errors": [],
            "warnings": [],
            "unused_variables": [],
            "missing_placeholders": []
        }
        
        try:
            prs = Presentation(self.template_path)
            
            # テンプレート内のプレースホルダーを収集
            template_placeholders = set()
            for slide in prs.slides:
                for slide_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        if shape.text_frame.text:
                            import re
                            placeholders = re.findall(r'\{\{[^}]+\}\}', shape.text_frame.text)
                            if placeholders:
                                template_placeholders.update(placeholders)
            
            # 提供された変数をチェック
            provided_variables = set(variables.keys())
            
            # 未使用の変数
            unused = provided_variables - template_placeholders
            if unused:
                validation_result["warnings"].append(f"未使用の変数: {list(unused)}")
                validation_result["unused_variables"] = list(unused)
            
            # 不足しているプレースホルダー
            missing = template_placeholders - provided_variables
            if missing:
                validation_result["errors"].append(f"不足している変数: {list(missing)}")
                validation_result["missing_placeholders"] = list(missing)
                validation_result["valid"] = False
            
            # 空の値のチェック
            empty_values = [k for k, v in variables.items() if not v or v.strip() == ""]
            if empty_values:
                validation_result["warnings"].append(f"空の値を持つ変数: {empty_values}")
            
        except Exception as e:
            validation_result["valid"] = False
            validation_result["errors"].append(f"テンプレート読み込みエラー: {e!s}")
        
        return validation_result


def create_temp_template(template_path: str, output_dir: str = None) -> str:
    """
    テンプレートの一時コピーを作成
    
    Args:
        template_path: 元のテンプレートパス
        output_dir: 出力ディレクトリ（Noneの場合は一時ディレクトリ）
        
    Returns:
        一時テンプレートのパス
    """
    import tempfile
    
    if output_dir is None:
        output_dir = tempfile.mkdtemp()
    
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    
    template_path = Path(template_path)
    temp_template = output_dir / f"temp_{template_path.name}"
    
    shutil.copy2(template_path, temp_template)
    
    return str(temp_template)


def cleanup_temp_template(temp_path: str):
    """一時テンプレートを削除"""
    try:
        os.remove(temp_path)
        print(f"一時テンプレートを削除しました: {temp_path}")
    except Exception as e:
        print(f"一時テンプレート削除エラー: {e}")
