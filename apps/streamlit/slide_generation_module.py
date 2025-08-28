"""
slide_generation_module.py
このモジュールは、営業提案スライドの生成に使用されるStreamlit UIと関連ビジネスロジックを定義します。
UIは元のアプリケーションのレイアウト（グローバル設定用のサイドバーと、製品推奨とスライド生成のための2段階ワークフロー）を踏襲していますが、
基盤となる推奨ロジックは以前の動作実装から復元されています。復元されたロジックは、言語モデルを使用して会議メモ、アップロードコンテンツ、チャット履歴を分析して問題点を特定し、キーワードマッチングとベクトル類似度検索の両方を使用してCSVカタログから候補製品をランク付けし、最後にLLMを使用して上位製品を選択して要約します。
LLM呼び出しが失敗した場合、システムはより基本的なヒューリスティックにフォールバックします。
Streamlitはこのファイルをアプリ内で実行するため、ほとんどの関数は純粋で副作用がありません。
セッション状態は、ユーザー入力と中間結果（候補製品のリストや抽出された問題など）を再実行間で保持するために使用されます。 
UI の定義については、このモジュールの下部にある `render_slide_generation_page` 関数を参照してください。
"""

from __future__ import annotations

import json
import os
import re
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any, List, Dict

import numpy as np
import pandas as pd
from dotenv import load_dotenv

# Make sure environment variables (API keys, etc.) are loaded.  This
# mirrors behaviour from the original code base.
load_dotenv(".env", override=True)

import streamlit as st

# API client used for retrieving chat history.  This mirrors the
# behaviour in the existing codebase where an internal API returns
# previous messages for a given project.
from lib.api import api_available, get_api_client

# Styling helpers and the slide generator class.  These utilities are
# unchanged from the original implementation and are imported here for
# clarity.
from lib.styles import (
    apply_company_analysis_page_styles,
    apply_main_styles,
    apply_slide_generation_page_styles,
    apply_title_styles,
    render_sidebar_logo_card,
    render_slide_generation_title,
)
from lib.new_slide_generator import NewSlideGenerator

# Attempt to import the OpenAI clients.  These are optional and will be
# used only if API keys are configured.  The restored logic uses the
# same approach as the working version: it will automatically select
# between Azure and OpenAI depending on environment variables.
try:
    from openai import AzureOpenAI, OpenAI
except Exception:
    AzureOpenAI = None  # type: ignore
    OpenAI = None  # type: ignore


###############################################################################
# Constants and configuration
###############################################################################

# Determine the project root relative to this file.  Assets such as the logo
# and product CSVs are located relative to this root.
PROJECT_ROOT = Path(__file__).parent.parent.parent
LOGO_PATH = PROJECT_ROOT / "data" / "images" / "otsuka_logo.jpg"
ICON_PATH = PROJECT_ROOT / "data" / "images" / "otsuka_icon.png"
PRODUCTS_DIR = PROJECT_ROOT / "data" / "csv" / "products"
PLACEHOLDER_IMG = PROJECT_ROOT / "data" / "images" / "product_placeholder.png"

# Path to the SQLite database used for persisting proposal drafts.  It
# resides under `data/sqlite` at the project root.  The helper
# functions `_init_db_for_proposals` and `_save_proposal_to_db` rely on
# this path.
DB_PATH = PROJECT_ROOT / "data" / "sqlite" / "app.db"


###############################################################################
# Session initialisation
###############################################################################

def _ensure_session_defaults() -> None:
    """Populate session state with default values."""
    ss = st.session_state
    # Context from the selected project (set elsewhere in the app)
    ss.setdefault("selected_project", None)
    ss.setdefault("api_error", None)
    # User input: meeting notes and uploaded files
    ss.setdefault("slide_meeting_notes", "")
    ss.setdefault("uploaded_files_store", [])  # file_uploader state
    # Result lists
    ss.setdefault("product_candidates", [])  # candidate products (list of dicts)
    ss.setdefault("analyzed_issues", [])  # extracted pain points
    # Draft outline and overview input for slide generation
    ss.setdefault("slide_outline", None)
    ss.setdefault("slide_overview", "")
    # Configuration options; defaults mirror the original UI
    ss.setdefault("slide_history_reference_count", 3)
    ss.setdefault("slide_top_k", 1)
    ss.setdefault("slide_products_dataset", "Auto")
    ss.setdefault("slide_use_tavily_api", True)
    ss.setdefault("slide_use_gpt_api", True)
    ss.setdefault("slide_tavily_uses", 1)
    # Embedding cache used for similarity search
    ss.setdefault("_emb_cache", {})
    # Template upload bytes and name used for slide generation
    ss.setdefault("slide_template_bytes", None)
    ss.setdefault("slide_template_name", None)
    # Last saved proposal ID for retrieving issues later when generating slides
    ss.setdefault("last_proposal_id", None)


###############################################################################
# Utility functions
###############################################################################

def _get_proposal_issues_from_db(proposal_id: str) -> List[Dict[str, Any]]:
    """Retrieve previously saved proposal issues from the SQLite database."""
    if not proposal_id:
        return []
    try:
        import sqlite3
        if not DB_PATH.exists():
            return []
        with sqlite3.connect(DB_PATH) as conn:
            cursor = conn.cursor()
            cursor.execute(
                """
                SELECT idx, issue, weight, keywords_json
                FROM proposal_issues
                WHERE proposal_id = ?
                ORDER BY idx
                """,
                (proposal_id,),
            )
            rows = cursor.fetchall()
            issues: List[Dict[str, Any]] = []
            for row in rows:
                import json as _j
                keywords = _j.loads(row[3]) if row[3] else []
                issues.append(
                    {"issue": row[1], "weight": row[2], "keywords": keywords}
                )
            return issues
    except Exception:
        return []


def _to_float(val) -> float | None:
    """Convert a value representing a price to a float."""
    if val is None:
        return None
    if isinstance(val, (int, float)):
        if isinstance(val, float) and pd.isna(val):
            return None
        return float(val)
    if isinstance(val, str):
        s = val.strip().replace("¥", "").replace(",", "")
        if not s:
            return None
        try:
            return float(s)
        except Exception:
            return None
    return None


def _fmt_price(val) -> str:
    """Format a numerical price value as a string with a yen symbol and commas."""
    v = _to_float(val)
    return f"¥{int(round(v)):,}" if v is not None else "—"


def _list_product_datasets() -> List[str]:
    """Return a list of subfolder names under PRODUCTS_DIR."""
    if not PRODUCTS_DIR.exists():
        return ["Auto"]
    ds = ["Auto"]
    for p in PRODUCTS_DIR.iterdir():
        if p.is_dir():
            ds.append(p.name)
    return ds


def _gather_messages_context(item_id: str | None, history_n: int) -> str:
    """Return the last N message exchanges (2N messages) for a project item."""
    if not (item_id and api_available()):
        return ""
    try:
        api = get_api_client()
        msgs = api.get_item_messages(item_id) or []
        take = min(len(msgs), history_n * 2)
        recent = msgs[-take:] if take > 0 else []
        ctx_lines = []
        for m in recent:
            role = m.get("role", "assistant")
            role_j = "ユーザー" if role == "user" else "アシスタント"
            ctx_lines.append(f"{role_j}: {m.get('content','')}")
        return "\n".join(ctx_lines)
    except Exception:
        return ""


def _load_products_from_csv(dataset: str) -> pd.DataFrame:
    """Load product catalogues from CSV files."""
    frames: list[pd.DataFrame] = []
    if not PRODUCTS_DIR.exists():
        return pd.DataFrame()

    def _read_csvs(folder: Path) -> None:
        for csvp in folder.glob("*.csv"):
            try:
                df = pd.read_csv(csvp)
                # Ensure expected columns exist
                for col in ["name", "category", "price", "description", "tags"]:
                    if col not in df.columns:
                        df[col] = None
                for col in ["image_url", "image", "thumbnail"]:
                    if col not in df.columns:
                        df[col] = None
                if "id" not in df.columns:
                    df["id"] = [f"{csvp.stem}-{i+1}" for i in range(len(df))]
                df["source_csv"] = csvp.stem
                frames.append(
                    df[[
                        "id",
                        "name",
                        "category",
                        "price",
                        "description",
                        "tags",
                        "image_url",
                        "image",
                        "thumbnail",
                        "source_csv",
                    ]]
                )
            except Exception:
                continue

    if dataset == "Auto":
        for sub in PRODUCTS_DIR.iterdir():
            if sub.is_dir():
                _read_csvs(sub)
        _read_csvs(PRODUCTS_DIR)
    else:
        target = PRODUCTS_DIR / dataset
        if target.exists() and target.is_dir():
            _read_csvs(target)
    if not frames:
        return pd.DataFrame()
    return pd.concat(frames, ignore_index=True)


def _extract_text_from_uploads(
    uploaded_files: List[Any], max_chars: int = 12000
) -> str:
    """Extract and concatenate text from a list of uploaded files."""
    if not uploaded_files:
        return ""
    chunks: list[str] = []
    used_chars = 0

    def _append(text: str) -> None:
        nonlocal used_chars
        if not text:
            return
        remain = max_chars - used_chars
        if remain <= 0:
            return
        cut = text[:remain]
        chunks.append(cut)
        used_chars += len(cut)

    for f in uploaded_files:
        try:
            name = getattr(f, "name", "uploaded_file")
            lower = str(name).lower()
            # Reset file pointer
            try:
                f.seek(0)
            except Exception:
                pass
            data = f.read()
            try:
                f.seek(0)
            except Exception:
                pass
            if lower.endswith(".pdf"):
                try:
                    import io
                    from pypdf import PdfReader
                    reader = PdfReader(io.BytesIO(data))
                    page_limit = min(len(reader.pages), 30)
                    texts: list[str] = []
                    for i in range(page_limit):
                        try:
                            texts.append(reader.pages[i].extract_text() or "")
                        except Exception:
                            continue
                    _append(f"\n[PDF:{name} 抜粋]\n" + "\n".join(texts))
                except Exception:
                    _append(f"\n[PDF:{name}]（抽出失敗→ファイル名のみ反映）")
            elif lower.endswith(".docx"):
                try:
                    import io
                    from docx import Document
                    doc = Document(io.BytesIO(data))
                    paras = [p.text for p in doc.paragraphs if p.text]
                    _append(f"\n[DOCX:{name} 抜粋]\n" + "\n".join(paras))
                except Exception:
                    _append(f"\n[DOCX:{name}]（抽出失敗→ファイル名のみ反映）")
            elif lower.endswith(".pptx"):
                try:
                    import io
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(data))
                    slide_texts: list[str] = []
                    for s in prs.slides:
                        buf: list[str] = []
                        for shp in s.shapes:
                            try:
                                if hasattr(shp, "text"):
                                    t = shp.text or ""
                                    if t:
                                        buf.append(t)
                            except Exception:
                                continue
                        if buf:
                            slide_texts.append("\n".join(buf))
                    _append(f"\n[PPTX:{name} 抜粋]\n" + "\n---\n".join(slide_texts))
                except Exception:
                    _append(f"\n[PPTX:{name}]（抽出失敗→ファイル名のみ反映）")
            elif lower.endswith(".csv"):
                try:
                    import io
                    tmp = io.BytesIO(data)
                    df = pd.read_csv(tmp)
                    head = df.head(20)
                    txt = head.to_csv(index=False)
                    _append(f"\n[CSV:{name} 先頭20行]\n{txt}")
                except Exception:
                    _append(f"\n[CSV:{name}]（抽出失敗→ファイル名のみ反映）")
            elif lower.endswith(".txt"):
                try:
                    txt = data.decode("utf-8", errors="ignore")
                except Exception:
                    txt = str(data[:4000])
                _append(f"\n[TXT:{name}]\n{txt}")
            else:
                _append(f"\n[{name}]（未対応/バイナリのため概要反映のみ）")
        except Exception:
            continue
        if used_chars >= max_chars:
            break
    return "\n".join(chunks).strip()


def _simple_tokenize(text: str) -> List[str]:
    """A simple tokenizer used for keyword matching."""
    text = str(text or "").lower()
    text = re.sub(r"[^a-z0-9\u3040-\u30ff\u4e00-\u9fff]+", " ", text)
    toks = text.split()
    return [t for t in toks if len(t) >= 2]


def _fallback_rank_products(
    notes: str,
    messages_ctx: str,
    products_df: pd.DataFrame,
    top_pool: int,
) -> List[Dict[str, Any]]:
    """Simple keyword based ranking of products."""
    if products_df.empty:
        return []
    query_text = (notes or "") + "\n" + (messages_ctx or "")
    q_tokens = _simple_tokenize(query_text)
    def _row_text(row: pd.Series) -> str:
        return " ".join([
            str(row.get("name") or ""),
            str(row.get("category") or ""),
            str(row.get("description") or ""),
            str(row.get("tags") or ""),
        ]).lower()
    scored: list[tuple[float, Dict[str, Any]]] = []
    for _, row in products_df.iterrows():
        t = _row_text(row)
        score = 0.0
        for tok in q_tokens:
            if tok in t:
                score += 1.0
        reason = f"一致語句数={int(score)}" if score > 0 else "一致なし（低スコア）"
        scored.append(
            (
                score,
                {
                    "id": row.get("id"),
                    "name": row.get("name"),
                    "category": row.get("category"),
                    "price": row.get("price"),
                    "description": row.get("description"),
                    "tags": row.get("tags"),
                    "image_url": row.get("image_url"),
                    "image": row.get("image"),
                    "thumbnail": row.get("thumbnail"),
                    "source_csv": row.get("source_csv"),
                    "score": round(float(score), 2),
                    "reason": reason,
                },
            )
        )
    scored.sort(key=lambda x: (x[0], str(x[1]["name"]).lower()), reverse=True)
    return [d for _, d in scored[:top_pool]]


###############################################################################
# LLM integration helpers
###############################################################################

def _get_chat_client():
    """Return an OpenAI or Azure OpenAI chat client and the model name."""
    use_azure = os.getenv("USE_AZURE", "").lower() == "true" or bool(os.getenv("AZURE_OPENAI_ENDPOINT"))
    if use_azure:
        endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
        api_key = os.getenv("AZURE_OPENAI_API_KEY")
        api_version = os.getenv("API_VERSION", "2024-06-01")
        deployment = os.getenv("AZURE_OPENAI_CHAT_DEPLOYMENT")
        if not (endpoint and api_key and deployment):
            raise RuntimeError("Azure設定不足: AZURE_OPENAI_ENDPOINT / AZURE_OPENAI_API_KEY / AZURE_OPENAI_CHAT_DEPLOYMENT")
        client = AzureOpenAI(azure_endpoint=endpoint, api_key=api_key, api_version=api_version)
        model = deployment
    else:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise RuntimeError("OPENAI_API_KEY が未設定です。")
        client = OpenAI(api_key=api_key)
        model = os.getenv("DEFAULT_MODEL", "gpt-4o-mini")
    return client, model


def _normalize_concat_row(row: pd.Series) -> str:
    """Concatenate and normalise fields for embedding."""
    s = f"{row.get('name','')} {row.get('category','')} {row.get('tags','')} {row.get('description','')}"
    s = s.lower()
    s = re.sub(r"\s+", " ", s).strip()
    return s


def _embed_texts(client, texts: List[str], embed_model: str, is_azure: bool) -> np.ndarray:
    """Call the embedding API and return vectors."""
    try:
        resp = client.embeddings.create(model=embed_model, input=texts)
        vecs = np.array([d.embedding for d in resp.data], dtype="float32")
        return vecs
    except Exception as e:
        raise RuntimeError(f"embedding failed: {e}")


def _build_products_index(
    dataset: str, df: pd.DataFrame, client, embed_model: str, is_azure: bool
) -> Dict[str, Any]:
    """Construct an embedding index over the product catalogue."""
    key = f"{dataset}:{len(df)}:{embed_model}"
    cache = st.session_state.get("_emb_cache", {})
    if key in cache:
        return cache[key]
    texts = [_normalize_concat_row(row) for _, row in df.iterrows()]
    try:
        vecs = _embed_texts(client, texts, embed_model, is_azure)
        index = {
            "vecs": vecs,
            "ids": df["id"].astype(str).tolist(),
            "df": df,
            "model": embed_model,
        }
    except Exception:
        # Fallback: build a TF-IDF vectoriser
        try:
            from sklearn.feature_extraction.text import TfidfVectorizer
            vectorizer = TfidfVectorizer(min_df=1)
            vecs = vectorizer.fit_transform(texts)
            index = {
                "vecs": vecs,
                "ids": df["id"].astype(str).tolist(),
                "df": df,
                "model": "tfidf",
                "vectorizer": vectorizer,
            }
        except Exception:
            index = {"vecs": None, "ids": [], "df": df, "model": None}
    cache[key] = index
    st.session_state._emb_cache = cache
    return index


def _retrieve_by_issues(
    index: Dict[str, Any],
    issues: List[Dict[str, Any]],
    client,
    embed_model: str,
    is_azure: bool,
    top_pool: int,
) -> List[Dict[str, Any]]:
    """Perform weighted similarity search against the product index."""
    if not issues or not index or index.get("vecs") is None:
        return []
    vecs = index["vecs"]
    # If TF-IDF fallback is used, we skip similarity search
    if hasattr(vecs, "toarray") or index.get("model") == "tfidf":
        return []
    # Construct weighted query vector
    queries = [f"{it['issue']} {' '.join(it.get('keywords') or [])}".strip() for it in issues]
    weights = np.array([float(it.get("weight", 0.0)) for it in issues], dtype="float32")
    try:
        q_embs = _embed_texts(client, queries, embed_model, is_azure)
        q = (weights[:, None] * q_embs).sum(axis=0, keepdims=True)
        # Normalise vectors
        v_norm = vecs / (np.linalg.norm(vecs, axis=1, keepdims=True) + 1e-9)
        q_norm = q / (np.linalg.norm(q, axis=1, keepdims=True) + 1e-9)
        sims = np.dot(q_norm, v_norm.T).ravel()
        order = sims.argsort()[::-1][: max(1, top_pool)]
        out: List[Dict[str, Any]] = []
        for idx_pos in order:
            rid = index["ids"][idx_pos]
            row = index["df"].iloc[idx_pos].to_dict()
            out.append(
                {
                    "id": row.get("id"),
                    "name": row.get("name"),
                    "category": row.get("category"),
                    "price": row.get("price"),
                    "description": row.get("description"),
                    "tags": row.get("tags"),
                    "image_url": row.get("image_url"),
                    "image": row.get("image"),
                    "thumbnail": row.get("thumbnail"),
                    "source_csv": row.get("source_csv"),
                    "score": float(sims[idx_pos]),
                    "reason": f"課題と高類似 ({sims[idx_pos]:.3f})",
                }
            )
        return out
    except Exception:
        return []


# -------------------- 変更点1: JSON抽出を強化 --------------------
def _extract_json(s: str) -> Dict[str, Any]:
    """
    - ```json ... ``` フェンス対応
    - 先頭/末尾に説明文があっても抽出
    - トップレベルが配列でも {"items":[...]} に包んで返す
    """
    s = (s or "").strip()
    if not s:
        return {}

    # ```json フェンス除去（```json / ``` どちらも許容）
    if s.startswith("```"):
        s = re.sub(r"^```(?:json)?\s*", "", s, flags=re.IGNORECASE)
        s = re.sub(r"\s*```$", "", s)

    # まずは素直に
    try:
        data = json.loads(s)
        if isinstance(data, list):
            return {"items": data}
        return data if isinstance(data, dict) else {}
    except Exception:
        pass

    # テキスト中の最初の { ... } を抜き出す
    try:
        start = s.find("{")
        end = s.rfind("}")
        if start >= 0 and end > start:
            data = json.loads(s[start:end + 1])
            if isinstance(data, list):
                return {"items": data}
            return data if isinstance(data, dict) else {}
    except Exception:
        return {}
    return {}


# -------------------- 変更点2: 安全な LLM 呼び出しヘルパ --------------------
# 置き換え：_safe_chat_json
def _safe_chat_json(messages: List[Dict[str, str]], *, require_json: bool = True, temperature: float = 0.2) -> Dict[str, Any]:
    """
    LLM呼び出し（Azure/OpenAI両対応）
    - 一部モデルが temperature をサポートしない → 自動で温度なしリトライ
    - 一部モデルが response_format=json をサポートしない → プレーン出力でリトライ
    - 失敗理由は st.session_state.api_error に格納
    """
    try:
        client, model = _get_chat_client()
    except Exception as e:
        st.session_state.api_error = f"LLMクライアント初期化に失敗: {e}"
        return {}

    def _attempt(pass_temperature: bool, pass_json_mode: bool):
        kwargs = {"model": model, "messages": messages}
        # temperature は「明示的に許される場合のみ」付与したいが、
        # 互換性のため最初の試行では付与 → 失敗時に温度なしで再試行する。
        if pass_temperature:
            kwargs["temperature"] = temperature
        if pass_json_mode and require_json:
            kwargs["response_format"] = {"type": "json_object"}
        return client.chat.completions.create(**kwargs)

    # 1) 温度あり + JSONモード → 2) 温度なし + JSON → 3) 温度なし + プレーン
    resp = None
    try:
        resp = _attempt(pass_temperature=True, pass_json_mode=True)
    except Exception as e1:
        msg1 = str(e1)
        # temperature 非対応や JSONモード非対応の可能性 → 温度なしで再試行
        try:
            resp = _attempt(pass_temperature=False, pass_json_mode=True)
        except Exception as e2:
            msg2 = str(e2)
            # さらに JSON モードも外して再試行（プレーンテキストからJSON抽出）
            try:
                resp = _attempt(pass_temperature=False, pass_json_mode=False)
            except Exception as e3:
                st.session_state.api_error = f"LLM呼び出しに失敗: {e3}"
                # ここまで来たら完全失敗
                return {}

    txt = (resp.choices[0].message.content or "").strip()
    data = _extract_json(txt)
    if not data and require_json:
        # 念のためプレーンでもう一度（既に試しているが、明示的再試行）
        try:
            resp2 = client.chat.completions.create(model=model, messages=messages)
            txt2 = (resp2.choices[0].message.content or "").strip()
            data = _extract_json(txt2)
        except Exception as e:
            st.session_state.api_error = f"LLM呼び出しに失敗: {e}"
            return {}

    return data if isinstance(data, dict) else {}



def _llm_pick_products(
    pool: List[Dict[str, Any]],
    top_k: int,
    company: str,
    notes: str,
    ctx: str,
    issues: List[Dict[str, Any]] | None = None,
) -> List[Dict[str, Any]]:
    """Select the top products from the pool using an LLM."""
    if not pool:
        return []

    # GPT未使用なら LLM選抜をスキップ（プール上位を採用）
    if not st.session_state.get("slide_use_gpt_api", True):
        return pool[:top_k]

    lines: List[str] = []
    for p in pool:
        desc = (p.get("description") or "")[:200]
        tags = (p.get("tags") or "")[:120]
        cat = p.get("source_csv") or p.get("category") or ""
        price = p.get("price")
        price_s = f"¥{int(_to_float(price)):,}" if _to_float(price) is not None else "—"
        lines.append(
            f"- id:{p['id']} | name:{p.get('name','')} | category:{cat} | price:{price_s} | tags:{tags} | desc:{desc}"
        )
    catalog = "\n".join(lines)
    issues_text = ""
    if issues:
        parts: List[str] = []
        for i, it in enumerate(issues):
            kw = ", ".join(it.get("keywords") or [])
            parts.append(f"[{i}] {it.get('issue')} (重み={it.get('weight'):.2f}; キーワード={kw})")
        issues_text = "\n".join(parts)

    # Define the JSON schema expected in the response
    if issues:
        schema = {
            "recommendations": [
                {
                    "id": "<id>",
                    "reason": "<120字以内>",
                    "confidence": 0.0,
                    "solved_issue_ids": [0],
                    "evidence": "<根拠抜粋>",
                }
            ]
        }
    else:
        schema = {
            "recommendations": [
                {
                    "id": "<id>",
                    "reason": "<120字以内>",
                    "confidence": 0.0,
                }
            ]
        }

    user = f"""あなたはB2Bプリセールスの提案プランナーです。
以下の会社情報と商談詳細、会話文脈に基づいて、候補カタログから Top-{top_k} の製品を選び、日本語で短い理由（120字以内）と信頼度(0-1)を付けてください。
必ずカタログに存在する id のみを使用してください。
出力は JSON のみで、以下のスキーマに従ってください: {json.dumps(schema, ensure_ascii=False)}
# 会社: {company or "(なし)"}
# 商談詳細: {notes or "(なし)"}
# 会話文脈: {ctx or "(なし)"}
# 課題一覧: {issues_text or "(なし)"}
# 候補カタログ: {catalog}
"""

    data = _safe_chat_json(
        [
            {"role": "system", "content": "あなたは正確で簡潔な日本語で回答するアシスタントです。"},
            {"role": "user", "content": user},
        ],
        require_json=True,
        temperature=0.1,
    )
    recs = data.get("recommendations", []) if isinstance(data, dict) else []
    if not recs:
        return []

    pool_map = {str(p["id"]): p for p in pool}
    out: List[Dict[str, Any]] = []
    for r in recs:
        pid = str(r.get("id", "")).strip()
        if not pid or pid not in pool_map:
            continue
        src = pool_map[pid]
        reason = (r.get("reason") or "").strip() or src.get("reason")
        conf = float(r.get("confidence", 0.0)) if r.get("confidence") is not None else float(src.get("score", 0.0))
        solved_ids = r.get("solved_issue_ids") if isinstance(r.get("solved_issue_ids"), list) else []
        evidence = (r.get("evidence") or "").strip()
        out.append(
            {
                **src,
                "reason": reason,
                "score": conf,
                "solved_issue_ids": solved_ids,
                "evidence": evidence,
            }
        )
        if len(out) >= top_k:
            break
    return out


# -------------------- 変更点4: 要約もチェックボックス反映＆安全化 --------------------
def _summarize_overviews_llm(cands: List[Dict[str, Any]]) -> None:
    """Summarise product descriptions into 80 Japanese characters using an LLM."""
    items: List[Dict[str, str]] = []
    has_any = False
    for c in cands:
        mat = c.get("description") or c.get("tags") or c.get("name") or ""
        if mat:
            has_any = True
            items.append({"id": str(c.get("id") or ""), "name": c.get("name") or "", "material": str(mat)[:600]})
    if not has_any:
        for c in cands:
            c["overview"] = "—"
        return

    # GPT未使用ならフォールバック
    if not st.session_state.get("slide_use_gpt_api", True):
        for c in cands:
            base = c.get("description") or c.get("tags") or ""
            c["overview"] = (base[:80] + ("…" if base and len(base) > 80 else "")) if base else "—"
        return

    payload = "\n".join([f"- id:{it['id']} / 名称:{it['name']}\n 内容:{it['material']}" for it in items])
    prompt = (
        "各製品の「製品概要」を日本語で1〜2文、最大80字で要約してください。事実の追加・誇張は禁止。\n"
        "出力は JSON のみ: {\"summaries\":[{\"id\":\"<id>\",\"overview\":\"<80字以内>\"}]}\n"
        "入力:\n" + payload
    )
    data = _safe_chat_json(
        [
            {"role": "system", "content": "あなたは簡潔で正確な日本語の要約を作るアシスタントです。"},
            {"role": "user", "content": prompt},
        ],
        require_json=True,
        temperature=0.2,
    )

    mp: Dict[str, str] = {}
    if isinstance(data, dict):
        for s in (data.get("summaries") or []):
            pid = str(s.get("id") or "")
            ov = (s.get("overview") or "").strip()
            if pid and ov:
                mp[pid] = ov

    for c in cands:
        pid = str(c.get("id") or "")
        base = c.get("description") or c.get("tags") or ""
        fallback = (base[:80] + ("…" if base and len(base) > 80 else "")) if base else "—"
        c["overview"] = mp.get(pid, fallback)


def _resolve_product_image_src(rec: Dict[str, Any]) -> str | None:
    """Resolve the best available image source for a product record."""
    for key in ("image_url", "image", "thumbnail"):
        v = rec.get(key)
        if not v:
            continue
        s = str(v).strip()
        if s.startswith("http://") or s.startswith("https://"):
            return s
        p = (PROJECT_ROOT / s).resolve()
        if p.exists():
            return str(p)
    if PLACEHOLDER_IMG.exists():
        return str(PLACEHOLDER_IMG)
    return None


###############################################################################
# Business logic: pain point analysis and product search
###############################################################################

# -------------------- 変更点3: 課題抽出の堅牢化＆GPT使用フラグ対応 --------------------
def _analyze_pain_points(
    notes: str,
    messages_ctx: str,
    uploads_text: str = "",
) -> List[Dict[str, Any]]:
    """Analyse meeting notes, chat context and uploaded text to identify issues."""
    issues: List[Dict[str, Any]] = []

    # GPT未使用なら即フォールバックへ
    if not st.session_state.get("slide_use_gpt_api", True):
        pass
    else:
        sys = "あなたはB2B提案の課題分析アシスタントです。日本語でJSONのみ出力してください。"
        uploads_section = f"\n\n資料抜粋:\n{uploads_text}" if uploads_text else ""
        user = (
            "以下の情報から、解決したい課題を3〜6件抽出し、各課題に重み(0〜1)と関連キーワード(3〜6語)を付けてJSONで出力してください。\n"
            "- 課題は具体的に表現する\n"
            "- 重みは合計が約1になるよう相対調整\n"
            "- 引用可能なら資料由来の観点も反映（機密/個人情報は抽象化）\n"
            '出力スキーマ: {"issues":[{"issue":"<80字以内>","weight":0.0,"keywords":["k1","k2","k3"]}]}\n'
            f"商談メモ: {notes}\n"
            f"会話文脈: {messages_ctx}\n"
            f"{uploads_section}"
        )

        data = _safe_chat_json(
            [
                {"role": "system", "content": sys},
                {"role": "user", "content": user},
            ],
            require_json=True,
            temperature=0.2,
        )

        cand = (data.get("issues") if isinstance(data, dict) else None) or (data.get("items") if isinstance(data, dict) else None)
        if isinstance(cand, list):
            for it in cand[:6]:
                issue = str(it.get("issue") or "").strip()
                if not issue:
                    continue
                weight = float(it.get("weight", 0.0))
                keywords = [str(k).strip() for k in (it.get("keywords") or []) if str(k).strip()]
                issues.append({
                    "issue": issue[:80],
                    "weight": max(0.0, min(1.0, weight)),
                    "keywords": keywords[:6],
                })

    # フォールバック（または LLMゼロ件時）
    if not issues:
        issues = [
            {"issue": "情報共有の改善", "weight": 0.34, "keywords": ["ナレッジ共有", "コミュニケーション", "ドキュメント"]},
            {"issue": "コスト最適化", "weight": 0.33, "keywords": ["費用削減", "効率化", "自動化"]},
            {"issue": "セキュリティ強化", "weight": 0.33, "keywords": ["アクセス管理", "監査", "権限"]},
        ]

    # 正規化
    s = sum(x["weight"] for x in issues) or 1.0
    for x in issues:
        x["weight"] = float(x["weight"] / s)
    return issues


def _search_product_candidates(
    company: str,
    item_id: str | None,
    meeting_notes: str,
    top_k: int,
    history_n: int,
    dataset: str,
    uploaded_files: List[Any],
    issues_precomputed: List[Dict[str, Any]] | None = None,
) -> List[Dict[str, Any]]:
    """Search and select top product candidates for a proposal."""
    # Retrieve recent chat history
    ctx = _gather_messages_context(item_id, history_n)
    # Load product catalogue
    df = _load_products_from_csv(dataset)
    if df.empty:
        return []
    # Extract text from uploaded materials
    uploads_text = _extract_text_from_uploads(uploaded_files) if uploaded_files else ""
    # Determine or compute issues
    issues = issues_precomputed if issues_precomputed is not None else _analyze_pain_points(meeting_notes or "", ctx or "", uploads_text)
    # Embedding model and environment detection
    use_azure = os.getenv("USE_AZURE", "").lower() == "true" or bool(os.getenv("AZURE_OPENAI_ENDPOINT"))
    embed_model = os.getenv("AZURE_OPENAI_EMBED_DEPLOYMENT") if use_azure else os.getenv("EMBED_MODEL", "text-embedding-3-small")
    # Build or retrieve the embedding index
    try:
        client, _ = _get_chat_client()
    except Exception as e:
        st.session_state.api_error = f"埋め込み用クライアント取得に失敗: {e}"
        client = None
    index = _build_products_index(dataset, df, client, embed_model, use_azure)
    # Perform weighted similarity search based on issues
    top_pool = max(40, top_k * 4)
    pool = _retrieve_by_issues(index, issues, client, embed_model, use_azure, top_pool)
    # Fallback: keyword based ranking if no embedding results
    if not pool:
        pool = _fallback_rank_products(meeting_notes, ctx, df, top_pool=max(40, top_k * 3))
    # Select final candidates using LLM; fall back to top of pool
    try:
        selected = _llm_pick_products(pool, top_k, company, meeting_notes, ctx, issues)
        if not selected:
            selected = pool[:top_k]
    except Exception:
        selected = pool[:top_k]
    # Summarise product descriptions
    _summarize_overviews_llm(selected)
    return selected


###############################################################################
# Proposal persistence
###############################################################################

def _init_db_for_proposals() -> None:
    """Initialise the proposals database if not already present."""
    DB_PATH.parent.mkdir(parents=True, exist_ok=True)
    import sqlite3
    with sqlite3.connect(DB_PATH) as conn:
        c = conn.cursor()
        c.execute(
            """
            CREATE TABLE IF NOT EXISTS proposals(
                id TEXT PRIMARY KEY,
                project_item_id TEXT,
                company TEXT NOT NULL,
                meeting_notes TEXT,
                overview TEXT,
                created_at TEXT NOT NULL
            )
            """
        )
        c.execute(
            """
            CREATE TABLE IF NOT EXISTS proposal_issues(
                proposal_id TEXT NOT NULL,
                idx INTEGER NOT NULL,
                issue TEXT NOT NULL,
                weight REAL,
                keywords_json TEXT,
                FOREIGN KEY(proposal_id) REFERENCES proposals(id)
            )
            """
        )
        c.execute(
            """
            CREATE TABLE IF NOT EXISTS proposal_products(
                proposal_id TEXT NOT NULL,
                rank INTEGER NOT NULL,
                product_id TEXT,
                name TEXT,
                category TEXT,
                price TEXT,
                reason TEXT,
                overview TEXT,
                score REAL,
                source_csv TEXT,
                image_url TEXT,
                FOREIGN KEY(proposal_id) REFERENCES proposals(id)
            )
            """
        )
        conn.commit()


def _save_proposal_to_db(
    project_item_id: str | None,
    company: str,
    meeting_notes: str,
    overview: str,
    issues: List[Dict[str, Any]],
    products: List[Dict[str, Any]],
    created_at_iso: str,
) -> str:
    """Persist a proposal draft to the database and return its ID."""
    _init_db_for_proposals()
    import sqlite3, uuid, json as _json
    pid = str(uuid.uuid4())
    with sqlite3.connect(DB_PATH) as conn:
        c = conn.cursor()
        c.execute(
            "INSERT INTO proposals(id, project_item_id, company, meeting_notes, overview, created_at) VALUES(?,?,?,?,?,?)",
            (pid, project_item_id, company, meeting_notes, overview, created_at_iso),
        )
        for i, it in enumerate(issues or []):
            c.execute(
                "INSERT INTO proposal_issues(proposal_id, idx, issue, weight, keywords_json) VALUES(?,?,?,?,?)",
                (
                    pid,
                    i + 1,
                    it.get("issue", ""),
                    float(it.get("weight") or 0.0),
                    _json.dumps(it.get("keywords") or [], ensure_ascii=False),
                ),
            )
        for r, p in enumerate(products or []):
            c.execute(
                """
                INSERT INTO proposal_products(
                    proposal_id, rank, product_id, name, category, price, reason, overview, score, source_csv, image_url
                ) VALUES(?,?,?,?,?,?,?,?,?,?,?)
                """,
                (
                    pid,
                    r + 1,
                    str(p.get("id", "")) or None,
                    p.get("name", ""),
                    p.get("source_csv") or p.get("category", ""),
                    str(p.get("price", "")),
                    p.get("reason", ""),
                    p.get("overview", ""),
                    float(p.get("score") or 0.0),
                    p.get("source_csv") or "",
                    p.get("image_url") or "",
                ),
            )
        conn.commit()
    return pid


###############################################################################
# Rendering helpers
###############################################################################

def _render_issues_body(issues: List[Dict[str, Any]], body_ph: st.delta_generator.DeltaGenerator) -> None:
    """Render a list of issues into the provided Streamlit placeholder."""
    body_ph.empty()
    with body_ph.container():
        if not issues:
            st.caption("『商品提案を作成』を押すと、商談メモ・履歴・参考資料から課題を自動抽出して表示します。")
            return
        for i, it in enumerate(issues, start=1):
            with st.container(border=True):
                st.markdown(f"**{i}. {it.get('issue','—')}**")
                st.caption(f"重み: {it.get('weight',0):.2f}")
                kws = it.get("keywords") or []
                if kws:
                    st.markdown("関連キーワード: " + " / ".join(kws))


def _render_candidates_body(recs: List[Dict[str, Any]], body_ph: st.delta_generator.DeltaGenerator) -> None:
    """Render a list of product recommendation cards into the placeholder."""
    body_ph.empty()
    with body_ph.container():
        if not recs:
            st.info("提案候補がありません。『商品提案を作成』を押してください。")
            return
        for r in recs:
            pid = str(r.get("id") or "")
            name = str(r.get("name") or "")
            cat_src = r.get("source_csv") or r.get("category") or "—"
            price_s = _fmt_price(r.get("price"))
            reason = r.get("reason") or "—"
            overview = r.get("overview") or "—"
            with st.container(border=True):
                c1, c2 = st.columns([1, 3], gap="medium")
                with c1:
                    img_src = _resolve_product_image_src(r)
                    if img_src and (img_src.startswith("http") or os.path.exists(img_src)):
                        st.image(img_src, use_container_width=True)
                    else:
                        st.markdown(
                            "<div style='width:100%;height:120px;border:1px solid #eee;border-radius:10px;"
                            "background:#f6f7f9;display:flex;align-items:center;justify-content:center;"
                            "color:#999;font-size:12px;'>画像なし</div>",
                            unsafe_allow_html=True,
                        )
                with c2:
                    st.markdown(f"**{name}**")
                    st.caption(f"カテゴリ: {cat_src} ／ 価格: {price_s} ／ ID: {pid}")
                    st.markdown(f"**提案理由**：{reason}")
                    st.markdown(f"**製品概要**：{overview}")


###############################################################################
# Main rendering function
###############################################################################

def _make_outline_preview(company_name: str, meeting_notes: str, products: List[Dict[str, Any]], overview: str) -> Dict[str, Any]:
    """簡易アウトラインのプレビュー JSON を返す（既存互換のダミー実装）"""
    return {
        "company": company_name,
        "overview": overview,
        "meeting_notes": meeting_notes[:400],
        "products": [
            {"id": p.get("id"), "name": p.get("name"), "reason": p.get("reason"), "overview": p.get("overview")}
            for p in products
        ],
    }


def render_slide_generation_page() -> None:
    """Entry point to render the slide generation page in Streamlit."""
    _ensure_session_defaults()
    try:
        st.set_page_config(
            page_title="スライド作成",
            page_icon=str(ICON_PATH),
            layout="wide",
            initial_sidebar_state="expanded",
        )
    except Exception:
        # In case set_page_config is called outside of the main script
        pass
    # Apply common styles
    apply_main_styles(hide_sidebar=False, hide_header=True)
    apply_title_styles()
    apply_company_analysis_page_styles()
    apply_slide_generation_page_styles()
    # Determine project context (selected in a different part of the app)
    pj = st.session_state.get("selected_project")
    if pj:
        title_text = f"スライド作成 - {pj['title']} / {pj['company']}"
        company_internal = pj.get("company", "")
        item_id = pj.get("id")
    else:
        title_text = "スライド作成"
        company_internal = ""
        item_id = None
    # Sidebar with global settings
    with st.sidebar:
        render_sidebar_logo_card(LOGO_PATH)
        st.markdown("### 設定")
        st.text_input("企業名", value=company_internal, key="slide_company_input", disabled=True)
        st.selectbox(
            "商材データセット",
            options=_list_product_datasets(),
            key="slide_products_dataset",
            help="data/csv/products/ 配下のフォルダ。Autoは自動選択。",
        )
        st.markdown("---")
        st.markdown("### AI設定")
        st.checkbox("GPT API使用", key="slide_use_gpt_api")
        st.checkbox("TAVILY API使用", key="slide_use_tavily_api")
        if st.session_state.slide_use_tavily_api:
            st.selectbox(
                "TAVILY API呼び出し回数（製品あたり）",
                options=list(range(1, 6)),
                key="slide_tavily_uses",
            )
        sidebar_clear = st.button("クリア", use_container_width=True, help="提案候補と課題の表示をクリア")
        st.markdown("<div class='sidebar-bottom'>", unsafe_allow_html=True)
        if st.button("← 案件一覧に戻る", use_container_width=True):
            st.session_state.current_page = "案件一覧"
            st.session_state.page_changed = True
            st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)
    # Title header
    render_slide_generation_title(title_text)
    # Header row with step title and action button
    head_l, head_r = st.columns([8, 2])
    with head_l:
        st.subheader("1. 商品提案を作成")
    with head_r:
        search_btn = st.button("商品提案作成", type="primary", use_container_width=True)
    st.divider()
    # Top input area: meeting notes and file uploads
    top_l, top_r = st.columns([3, 2], gap="large")
    with top_l:
        st.markdown("**● 商談メモを入力**")
        st.text_area(
            label="商談メモを入力",
            key="slide_meeting_notes",
            height=154,
            label_visibility="collapsed",
            placeholder="例：来期の需要予測精度向上と在庫最適化。PoCから段階導入… など",
        )
    with top_r:
        st.markdown("**● 参考資料を入力（任意）**")
        uploads = st.file_uploader(
            label="参考資料を入力（任意）",
            type=["pdf", "pptx", "docx", "csv", "png", "jpg", "jpeg", "txt"],
            accept_multiple_files=True,
            key="slide_uploader",
            label_visibility="collapsed",
            help="議事録や要件定義などを添付。内容は課題抽出・候補選定に反映されます。",
        )
        if uploads:
            st.session_state.uploaded_files_store = uploads
            st.success(f"{len(uploads)} ファイルを受け付けました。")
        elif st.session_state.uploaded_files_store:
            st.caption(f"前回アップロード済み: {len(st.session_state.uploaded_files_store)} ファイル")
        # Additional detail settings in an expander
        with st.expander("詳細設定（商品提案の条件）", expanded=False):
            cols = st.columns(2)
            with cols[0]:
                st.selectbox(
                    "提案候補の件数（Top-K）",
                    options=list(range(1, 11)),
                    key="slide_top_k",
                    help="表示する提案候補の件数。",
                )
            with cols[1]:
                st.selectbox(
                    "過去ログ参照範囲（往復数）",
                    options=list(range(1, 11)),
                    key="slide_history_reference_count",
                    help="企業分析チャットの直近N往復を文脈として使用します。",
                )
    # Middle result area: issues and candidate lists
    bottom_l, bottom_r = st.columns([5, 7], gap="large")
    with bottom_l:
        st.markdown("**● 課題の要約（結果）**")
        issues_msg_ph = st.empty()
        issues_body_ph = st.empty()
    with bottom_r:
        st.markdown("**● 提案候補の一覧（結果）**")
        candidates_msg_ph = st.empty()
        candidates_body_ph = st.empty()
    # Render any previously stored results
    _render_issues_body(st.session_state.get("analyzed_issues") or [], issues_body_ph)
    _render_candidates_body(st.session_state.get("product_candidates") or [], candidates_body_ph)
    # Sidebar clear behaviour
    if sidebar_clear:
        st.session_state.product_candidates = []
        st.session_state.analyzed_issues = []
        st.session_state.slide_outline = None
        # Clear progress messages
        issues_msg_ph.empty()
        candidates_msg_ph.empty()
        # Render placeholders with no data
        _render_issues_body([], issues_body_ph)
        _render_candidates_body([], candidates_body_ph)
        st.success("提案候補と課題の表示をクリアしました。")
        st.stop()
    # When the user triggers product search
    if search_btn:
        if not company_internal.strip():
            st.error("企業が選択されていません。案件一覧から企業を選んでください。")
        else:
            with issues_msg_ph.container():
                with st.spinner("1/2 課題を抽出しています..."):
                    issues_body_ph.empty()
                    ctx_for_view = _gather_messages_context(item_id, int(st.session_state.slide_history_reference_count))
                    uploads_text_for_view = _extract_text_from_uploads(st.session_state.uploaded_files_store) if st.session_state.uploaded_files_store else ""
                    issues_early = _analyze_pain_points(
                        st.session_state.slide_meeting_notes or "",
                        ctx_for_view or "",
                        uploads_text_for_view or "",
                    )
                    st.session_state.analyzed_issues = issues_early
            # 課題抽出後の失敗理由をUI表示（変更点）
            issues_msg_ph.empty()
            if st.session_state.get("api_error"):
                st.warning(f"LLMの実行で問題が発生しました: {st.session_state.api_error}")
            _render_issues_body(issues_early, issues_body_ph)
            # Now select product candidates
            with candidates_msg_ph.container():
                with st.spinner("2/2 カタログと照合して提案候補を選定中..."):
                    candidates_body_ph.empty()
                    candidates = _search_product_candidates(
                        company=company_internal,
                        item_id=item_id,
                        meeting_notes=st.session_state.slide_meeting_notes or "",
                        top_k=int(st.session_state.slide_top_k),
                        history_n=int(st.session_state.slide_history_reference_count),
                        dataset=st.session_state.slide_products_dataset,
                        uploaded_files=st.session_state.uploaded_files_store,
                        issues_precomputed=issues_early,
                    )
                    st.session_state.product_candidates = candidates
            # Render candidates
            candidates_msg_ph.empty()
            _render_candidates_body(candidates, candidates_body_ph)
            # Save a draft of the proposal (issues and products) for later retrieval
            try:
                proposal_id = _save_proposal_to_db(
                    project_item_id=item_id,
                    company=company_internal,
                    meeting_notes=st.session_state.slide_meeting_notes or "",
                    overview=st.session_state.slide_overview or "",
                    issues=st.session_state.analyzed_issues or [],
                    products=st.session_state.product_candidates or [],
                    created_at_iso=datetime.now().isoformat(timespec="seconds"),
                )
                st.session_state["last_proposal_id"] = proposal_id
            except Exception as e:
                st.warning(f"ドラフト保存に失敗しました: {e}")
    # Second section: slide generation
    st.subheader("2. 提案スライド生成")
    st.divider()
    tmpl_file = st.file_uploader(
        "テンプレート（.pptx）を添付（任意）",
        type=["pptx"],
        key="slide_template_uploader",
        help="未添付の場合は既定テンプレートを使用します",
    )
    # Store uploaded template for later
    if tmpl_file is not None:
        st.session_state.slide_template_bytes = tmpl_file.getvalue()
        st.session_state.slide_template_name = tmpl_file.name
        st.success(f"テンプレートを受け付けました：{tmpl_file.name}")
    else:
        current = st.session_state.get("slide_template_name")
        if current:
            st.caption(f"現在のテンプレート：{current}（アップロード済みを使用）")
        else:
            st.caption("テンプレート未添付：既定テンプレートを使用します")
    # Overview input and generate button
    row_l, row_r = st.columns([8, 2], vertical_alignment="center")
    with row_l:
        st.session_state.slide_overview = st.text_input(
            "概説（任意）",
            value=st.session_state.slide_overview or "",
            placeholder="例：在庫最適化を中心に、需要予測と補充計画の連携を提案…",
            label_visibility="collapsed",
        )
    with row_r:
        gen_btn = st.button("生成", type="primary", use_container_width=True)
    if gen_btn:
        if not company_internal.strip():
            st.error("企業が選択されていません。")
        elif not st.session_state.product_candidates:
            st.error("提案候補がありません。先に『商品提案を作成』を押してください。")
        else:
            selected = list(st.session_state.product_candidates or [])  # adopt all candidates
            # Build draft outline for preview
            outline = _make_outline_preview(
                company_internal,
                st.session_state.slide_meeting_notes or "",
                selected,
                st.session_state.slide_overview or "",
            )
            st.session_state.slide_outline = outline
            # Prepare chat history
            chat_history = _gather_messages_context(item_id, st.session_state.slide_history_reference_count)
            # Prepare template file if uploaded
            uploaded_template_path: str | None = None
            if st.session_state.get("slide_template_bytes"):
                try:
                    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
                    tmp.write(st.session_state["slide_template_bytes"])
                    tmp.flush()
                    tmp.close()
                    uploaded_template_path = tmp.name
                except Exception:
                    uploaded_template_path = None
            # Generate presentation
            with st.spinner("AIエージェントがプレゼンテーションを生成中..."):
                try:
                    if uploaded_template_path:
                        generator = NewSlideGenerator(template_path=uploaded_template_path)
                    else:
                        generator = NewSlideGenerator()
                    pptx_data = generator.create_presentation(
                        project_name=company_internal,
                        company_name=company_internal,
                        meeting_notes=st.session_state.slide_meeting_notes or "",
                        chat_history=chat_history,
                        products=selected,
                        # ↓↓↓ 修正：未定義の proposal_issues を渡さない。DBから取得したものだけを渡す
                        proposal_issues=_get_proposal_issues_from_db(st.session_state.get("last_proposal_id") or ""),
                        proposal_id=st.session_state.get("last_proposal_id"),
                        use_tavily=st.session_state.slide_use_tavily_api,
                        use_gpt=st.session_state.slide_use_gpt_api,
                        tavily_uses=st.session_state.slide_tavily_uses,
                    )
                    # Present download button
                    st.success("プレゼンテーションが生成されました！")
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    filename = f"{company_internal}_提案書_{timestamp}.pptx"
                    st.download_button(
                        label="📥 プレゼンテーションをダウンロード",
                        data=pptx_data,
                        file_name=filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                        type="primary",
                    )
                except Exception as e:
                    st.error(f"プレゼンテーション生成でエラーが発生しました: {e}")
                    st.info("下書きのみ作成されました。")
                finally:
                    if uploaded_template_path and os.path.exists(uploaded_template_path):
                        try:
                            os.remove(uploaded_template_path)
                        except Exception:
                            pass
            # Show outline preview
            if st.session_state.slide_outline:
                with st.expander("下書きプレビュー（JSON）", expanded=True):
                    st.json(st.session_state.slide_outline)
