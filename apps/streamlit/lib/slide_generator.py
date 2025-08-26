"""
AI Agent for PPTX presentation generation
Creates presentations based on company reports, products and user input
"""

import os
import re
import json
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime

import streamlit as st
from openai import AzureOpenAI
from tavily import TavilyClient
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from dotenv import load_dotenv
from PIL import Image

# Загружаем переменные окружения
load_dotenv()

class SlideGenerator:
    """AI Agent for PPTX presentation generation"""
    
    def __init__(self):
        self.azure_client = None
        self.tavily_client = None
        self._init_clients()
    
    def _init_clients(self):
        """APIクライアントの初期化"""
        # Azure OpenAI
        if os.getenv("AZURE_OPENAI_API_KEY") and os.getenv("AZURE_OPENAI_ENDPOINT"):
            self.azure_client = AzureOpenAI(
                api_version="2024-12-01-preview",
                azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
                api_key=os.getenv("AZURE_OPENAI_API_KEY"),
            )
        
        # Tavily
        if os.getenv("TAVILY_API_KEY"):
            self.tavily_client = TavilyClient(api_key=os.getenv("TAVILY_API_KEY"))
    
    def generate_presentation(
        self,
        company_name: str,
        company_report: str,
        user_input: str,
        llm_proposal: str,
        additional_instructions: str,
        use_tavily_api: bool = False,
        use_gpt_api: bool = False,
        tavily_uses: int = 3
    ) -> bytes:
        """
        PPTXプレゼンテーションを生成
        
        Args:
            company_name: 企業名
            company_report: データベースからの企業レポート
            user_input: ユーザー入力（商談の詳細）
            llm_proposal: LLM提案（LLMによる提案）
            additional_instructions: 追加指示
            use_tavily_api: TAVILY APIを使用するか
            use_gpt_api: GPT APIを使用するか
            tavily_uses: 各商品に対するTAVILY API呼び出し回数
        
        Returns:
            bytes: バイト形式のPPTXファイル
        """
        try:
            # LLM提案をパースして商品を抽出
            products = self._parse_llm_proposal(llm_proposal)
            
            # APIが有効な場合、追加情報を取得
            if use_tavily_api and self.tavily_client:
                products = self._enhance_products_with_tavily(products, tavily_uses)
            else:
                # デフォルトテキストを使用
                for product in products:
                    product['tavily_info'] = "あいうえお"
            
            # GPTが有効な場合、スライド内容を生成
            if use_gpt_api and self.azure_client:
                slide_content = self._generate_slide_content_with_gpt(
                    company_name, company_report, user_input, products, additional_instructions
                )
            else:
                # デフォルト内容を使用
                slide_content = self._generate_default_slide_content(
                    company_name, company_report, user_input, products, additional_instructions
                )
            
            # プレゼンテーションを作成
            presentation = self._create_pptx_presentation(slide_content)
            
            # プレゼンテーションをバイト形式で返す（一時ファイルは作成しない）
            from io import BytesIO
            buffer = BytesIO()
            presentation.save(buffer)
            pptx_bytes = buffer.getvalue()
            buffer.close()
            
            return pptx_bytes
            
        except Exception as e:
            st.error(f"プレゼンテーション生成でエラーが発生しました: {str(e)}")
            raise
    
    def _parse_llm_proposal(self, llm_proposal: str) -> List[Dict[str, Any]]:
        """LLM提案をパースして商品を抽出"""
        products = []
        
        # 商品セクションを探す
        lines = llm_proposal.split('\n')
        in_products_section = False
        
        for line in lines:
            line = line.strip()
            
            # 商品セクションの開始
            if '推奨商材' in line or '推奨商材（max' in line:
                in_products_section = True
                continue
            
            # 商品セクションの終了
            if in_products_section and ('次のアクション' in line or '次アクション' in line):
                break
            
            # 商品をパース
            if in_products_section and line.startswith('- **'):
                product = self._parse_product_line(line)
                if product:
                    products.append(product)
        
        return products
    
    def _parse_product_line(self, line: str) -> Optional[Dict[str, Any]]:
        """商品行をパース"""
        try:
            # 例: - **Zowie CELERITAS**（カテゴリ：キーボード、概算価格：要見積／参考価格 ¥8,000〜¥15,000）
            # — 高耐久・高速入力が求められる営業・サポート部門の定番キーボード...
            
            # 商品名を抽出
            name_match = re.search(r'\*\*(.*?)\*\*', line)
            if not name_match:
                return None
            
            name = name_match.group(1)
            
            # カテゴリを抽出
            category_match = re.search(r'カテゴリ：([^、]+)', line)
            category = category_match.group(1) if category_match else "不明"
            
            # 価格を抽出
            price_match = re.search(r'参考価格\s*¥([0-9,]+)〜¥([0-9,]+)', line)
            if price_match:
                min_price = price_match.group(1).replace(',', '')
                max_price = price_match.group(2).replace(',', '')
                price = f"¥{min_price}〜¥{max_price}"
            else:
                price = "要見積"
            
            # 説明を抽出（—の後）
            description_match = re.search(r'—\s*(.+)', line)
            description = description_match.group(1) if description_match else ""
            
            return {
                'name': name,
                'category': category,
                'price': price,
                'description': description,
                'tavily_info': ""
            }
            
        except Exception as e:
            st.warning(f"商品パースでエラーが発生しました: {str(e)}")
            return None
    
    def _enhance_products_with_tavily(self, products: List[Dict[str, Any]], tavily_uses: int) -> List[Dict[str, Any]]:
        """TAVILY APIを通じて商品情報を強化"""
        if not self.tavily_client:
            return products
        
        for product in products:
            try:
                # 商品情報の検索
                query = f"{product['name']} {product['category']} 製品情報 特徴 仕様"
                search_result = self.tavily_client.search(
                    query=query,
                    search_depth="basic",
                    max_results=tavily_uses
                )
                
                # 関連情報を抽出
                relevant_info = []
                for result in search_result.get('results', [])[:tavily_uses]:
                    title = result.get('title', '')
                    content = result.get('content', '')
                    if title and content:
                        relevant_info.append(f"{title}: {content[:200]}...")
                
                product['tavily_info'] = "\n".join(relevant_info) if relevant_info else "あいうえお"
                
            except Exception as e:
                st.warning(f"{product['name']}のTAVILY検索でエラーが発生しました: {str(e)}")
                product['tavily_info'] = "あいうえお"
        
        return products
    
    def _generate_slide_content_with_gpt(
        self,
        company_name: str,
        company_report: str,
        user_input: str,
        products: List[Dict[str, Any]],
        additional_instructions: str
    ) -> Dict[str, Any]:
        """GPTを通じてスライド内容を生成"""
        if not self.azure_client:
            return self._generate_default_slide_content(
                company_name, company_report, user_input, products, additional_instructions
            )
        
        try:
            prompt = f"""
以下の情報を基に、B2B提案用のプレゼンテーションの内容を生成してください。

企業名: {company_name}

企業レポート:
{company_report}

商談の詳細:
{user_input}

提案商品:
{json.dumps(products, ensure_ascii=False, indent=2)}

追加指示:
{additional_instructions}

以下のスライドの内容を日本語で生成してください:

1. 現状の課題
- 企業が抱える問題点を分析
- 具体的な課題を3-5点挙げる

2. 導入メリット
- 導入による具体的な効果
- 定量的・定性的なメリット
- この内容は各商品の後に繰り返し表示されます

3. トータルコスト
- 総費用の概算
- 導入スケジュールとROI

注意: 商品の詳細は別々のスライドで自動生成され、各商品の後に導入メリットが表示されます。

出力形式: JSON形式で各スライドの内容を返してください。
各スライドは以下の形式で返してください:
{{
  "slide1": {{
    "title": "現状の課題",
    "content": ["内容1", "内容2", ...]
  }},
  "slide3": {{
    "title": "導入メリット",
    "content": ["内容1", "内容2", ...]
  }},
  "slide4": {{
    "title": "トータルコスト",
    "content": ["内容1", "内容2", ...]
  }}
}}
"""

            response = self.azure_client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "あなたはB2B提案の専門家です。具体的で実用的な提案内容を作成してください。"},
                    {"role": "user", "content": prompt}
                ],
                model="gpt-5-mini",
                max_tokens=4000,
                temperature=0.7
            )
            
            content = response.choices[0].message.content
            
            # JSONのパースを試行
            try:
                return json.loads(content)
            except json.JSONDecodeError:
                # JSONのパースに失敗した場合、デフォルト内容を使用
                st.warning("GPTが無効なJSONを返しました、デフォルト内容を使用します")
                return self._generate_default_slide_content(
                    company_name, company_report, user_input, products, additional_instructions
                )
                
        except Exception as e:
            st.warning(f"GPT生成でエラーが発生しました: {str(e)}, デフォルト内容を使用します")
            return self._generate_default_slide_content(
                company_name, company_report, user_input, products, additional_instructions
            )
    
    def _generate_default_slide_content(
        self,
        company_name: str,
        company_report: str,
        user_input: str,
        products: List[Dict[str, Any]],
        additional_instructions: str
    ) -> Dict[str, Any]:
        """デフォルトスライド内容を生成"""
        return {
            "slide1": {
                "title": "現状の課題",
                "content": [
                    f"{company_name}の現状分析",
                    "• 業務効率の低下",
                    "• システムの老朽化",
                    "• コストの増大",
                    "• セキュリティリスク",
                    "• 競合他社との差別化不足"
                ]
            },
            "slide2": {
                "title": "ご提案機器について",
                "content": [
                    "提案商品の詳細",
                    "",
                    *[f"• {p['name']} ({p['category']}) - {p['price']}" for p in products],
                    "",
                    "選定理由:",
                    *[f"• {p['description'][:100]}..." for p in products if p['description']]
                ],
                "products": products  # スライド用の商品情報を追加
            },
            "slide3": {
                "title": "導入メリット",
                "content": [
                    "導入による効果",
                    "• 業務効率の向上 (20-30%改善)",
                    "• コスト削減 (年間15-25%削減)",
                    "• セキュリティの強化",
                    "• ユーザー満足度の向上",
                    "• 競合優位性の確保"
                ]
            },
            "slide4": {
                "title": "トータルコスト",
                "content": [
                    "費用とスケジュール",
                    f"• 総投資額: 要見積 (概算: ¥{len(products) * 50000:,})",
                    "• 導入期間: 2-3ヶ月",
                    "• 期待ROI: 12-18ヶ月",
                    "• 運用コスト: 月額要見積",
                    "",
                    "次のステップ:",
                    "• 詳細見積の取得",
                    "• パイロット導入の検討",
                    "• 導入計画の策定"
                ]
            }
        }
    
    def _create_pptx_presentation(self, slide_content: Dict[str, Any]) -> Presentation:
        """PPTXプレゼンテーションを作成"""
        # 新しいプレゼンテーションを作成
        prs = Presentation()
        
        # スライドサイズを設定（16:9）
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # 最初のスライドを作成 - 現状の課題
        slide = self._create_slide(prs, "現状の課題", slide_content.get("slide1", {}).get("content", []))
        
        # 各商品のスライドを作成し、各商品の後に導入メリットを作成
        products = slide_content.get("slide2", {}).get("products", [])
        for i, product in enumerate(products):
            # 商品スライド
            slide_title = f"ご提案機器 {i+1}: {product['name']}"
            slide = self._create_single_product_slide(prs, slide_title, product)
            
            # 商品の後の導入メリットスライド
            benefits_content = slide_content.get("slide3", {}).get("content", [])
            slide = self._create_slide(prs, "導入メリット", benefits_content)
        
        # 最後のスライドを作成 - トータルコスト
        slide = self._create_slide(prs, "トータルコスト", slide_content.get("slide4", {}).get("content", []))
        
        # コスト比較スライドを作成
        slide = self._create_cost_comparison_slide(prs, products, slide_content)
        
        return prs
    
    def _create_slide(self, prs: Presentation, title: str, content: List[str]) -> Any:
        """個別のスライドを作成"""
        # タイトルとコンテンツのレイアウトを使用
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # 会社ロゴを追加
        self._add_logo_to_slide(slide)
        
        # スライドタイトル
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        # スライドコンテンツ
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, line in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 0, 0)
            
            # 見出しを太字で強調
            if line.startswith('•') or line.startswith('*'):
                p.font.bold = True
        
        return slide
    
    def _create_single_product_slide(self, prs: Presentation, title: str, product: Dict[str, Any]) -> Any:
        """単一商品のスライドを作成"""
        # 空白レイアウトを使用
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # 会社ロゴを追加
        self._add_logo_to_slide(slide)
        
        # スライドタイトル
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        # 商品情報（左側）
        y_position = 1.8
        
        # 商品名
        product_title = slide.shapes.add_textbox(Inches(0.5), Inches(y_position), Inches(8), Inches(0.8))
        title_frame = product_title.text_frame
        title_frame.text = f"商品名: {product['name']}"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        y_position += 1.0
        
        # カテゴリ
        category_text = slide.shapes.add_textbox(Inches(0.5), Inches(y_position), Inches(8), Inches(0.6))
        category_frame = category_text.text_frame
        category_frame.text = f"カテゴリ: {product['category']}"
        category_frame.paragraphs[0].font.size = Pt(18)
        category_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        y_position += 0.8
        
        # 価格
        price_text = slide.shapes.add_textbox(Inches(0.5), Inches(y_position), Inches(8), Inches(0.6))
        price_frame = price_text.text_frame
        price_frame.text = f"価格: {product['price']}"
        price_frame.paragraphs[0].font.size = Pt(18)
        price_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        y_position += 0.8
        
        # 説明
        if product.get('description'):
            desc_text = slide.shapes.add_textbox(Inches(0.5), Inches(y_position), Inches(8), Inches(2.0))
            desc_frame = desc_text.text_frame
            desc_frame.text = f"選定理由:\n{product['description']}"
            desc_frame.paragraphs[0].font.size = Pt(16)
            desc_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        # 商品画像（右側）
        try:
            # プロジェクトルートからの相対パスを使用
            project_root = Path(__file__).parent.parent.parent
            image_path = project_root / "data/images/example_picture.png"
            if image_path.exists():
                # 右側に画像を配置し、比例を保持
                img_left = Inches(9.5)
                img_top = Inches(1.8)
                img_width = Inches(3.5)
                
                # 元の画像サイズを取得して比例を保持
                with Image.open(image_path) as img:
                    orig_width, orig_height = img.size
                    aspect_ratio = orig_height / orig_width
                    img_height = img_width * aspect_ratio
                
                slide.shapes.add_picture(str(image_path), img_left, img_top, img_width, img_height)
        except Exception as e:
            # 画像の追加に失敗した場合、プレースホルダーを追加
            placeholder = slide.shapes.add_textbox(img_left, img_top, img_width, img_height)
            placeholder_frame = placeholder.text_frame
            placeholder_frame.text = "画像\nなし"
            placeholder_frame.paragraphs[0].font.size = Pt(20)
            placeholder_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
            placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_cost_comparison_slide(self, prs: Presentation, products: List[Dict[str, Any]], slide_content: Dict[str, Any]) -> Any:
        """コスト比較スライドを作成（LLMデータを使用）"""
        # 空白レイアウトを使用
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # 会社ロゴを追加
        self._add_logo_to_slide(slide)
        
        # タイトルバー（オレンジグラデーション風）
        title_bar = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        title_frame = title_bar.text_frame
        title_frame.text = "トータルコスト比較"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # タイトルバーの背景色を設定（オレンジ風）
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = RGBColor(255, 165, 0)
        
        # LLMデータからコスト情報を抽出
        cost_data = self._extract_cost_data_from_llm(products, slide_content)
        
        # 左列：現状
        current_section = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.5), Inches(5.5))
        current_frame = current_section.text_frame
        current_frame.text = f"現状\n{cost_data['current']['equipment']}\n\n【ハード料金】\n{cost_data['current']['hardware']}\n\n【ランニングコスト】\n{cost_data['current']['running']}\n\n【電気料金】\n{cost_data['current']['electricity']}\n※TEC値合計 {cost_data['current']['tec']}\n\n──────────────\n現状経費合計\n{cost_data['current']['total']}"
        
        # 現状セクションのスタイル設定
        for i, paragraph in enumerate(current_frame.paragraphs):
            if i == 0:  # タイトル
                paragraph.font.size = Pt(24)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            elif i == 1:  # サブタイトル
                paragraph.font.size = Pt(16)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            elif "【" in paragraph.text:  # カテゴリ
                paragraph.font.size = Pt(18)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(64, 64, 64)
            elif "円" in paragraph.text:  # 金額
                paragraph.font.size = Pt(20)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(64, 64, 64)
            elif "※" in paragraph.text:  # 注釈
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = RGBColor(128, 128, 128)
            elif "────────" in paragraph.text:  # 区切り線
                paragraph.font.size = Pt(16)
                paragraph.font.color.rgb = RGBColor(128, 128, 128)
            elif "現状経費合計" in paragraph.text:  # 合計ラベル
                paragraph.font.size = Pt(18)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(64, 64, 64)
            else:  # その他
                paragraph.font.size = Pt(16)
                paragraph.font.color.rgb = RGBColor(64, 64, 64)
        
        # 現状セクションの背景色（グレー）
        current_section.fill.solid()
        current_section.fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        # 右列：提案
        proposal_section = slide.shapes.add_textbox(Inches(6.5), Inches(1.8), Inches(5.5), Inches(5.5))
        proposal_frame = proposal_section.text_frame
        proposal_frame.text = f"提案\n{cost_data['proposal']['equipment']}\n\n【ハード料金】\n{cost_data['proposal']['hardware']}\n\n【ランニングコスト】\n{cost_data['proposal']['running']}\n\n【電気料金】\n{cost_data['proposal']['electricity']}\n※TEC値合計 {cost_data['proposal']['tec']}\n\n──────────────\nご提案経費合計\n{cost_data['proposal']['total']}"
        
        # 提案セクションのスタイル設定
        for i, paragraph in enumerate(proposal_frame.paragraphs):
            if i == 0:  # タイトル
                paragraph.font.size = Pt(24)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            elif i == 1:  # サブタイトル
                paragraph.font.size = Pt(16)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            elif "【" in paragraph.text:  # カテゴリ
                paragraph.font.size = Pt(18)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(64, 64, 64)
            elif "円" in paragraph.text:  # 金額
                paragraph.font.size = Pt(20)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(64, 64, 64)
            elif "※" in paragraph.text:  # 注釈
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = RGBColor(128, 128, 128)
            elif "────────" in paragraph.text:  # 区切り線
                paragraph.font.size = Pt(16)
                paragraph.font.color.rgb = RGBColor(128, 128, 128)
            elif "ご提案経費合計" in paragraph.text:  # 合計ラベル
                paragraph.font.size = Pt(18)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(64, 64, 64)
            else:  # その他
                paragraph.font.size = Pt(16)
                paragraph.font.color.rgb = RGBColor(64, 64, 64)
        
        # 提案セクションの背景色（青のヘッダー、白の本体）
        proposal_header = slide.shapes.add_textbox(Inches(6.5), Inches(1.8), Inches(5.5), Inches(0.8))
        proposal_header.fill.solid()
        proposal_header.fill.fore_color.rgb = RGBColor(0, 100, 200)
        
        # 下部サマリーエリア（赤いバー）
        summary_bar = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(11.5), Inches(1.0))
        summary_frame = summary_bar.text_frame
        summary_frame.text = f"月間経費差額 約 ▲{cost_data['difference']['monthly']}    年間経費差額 約 ▲{cost_data['difference']['yearly']}    5年間経費差額 約 ▲{cost_data['difference']['five_year']}"
        
        # サマリーのスタイル設定
        summary_frame.paragraphs[0].font.size = Pt(20)
        summary_frame.paragraphs[0].font.bold = True
        summary_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        summary_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # サマリーバーの背景色（赤）
        summary_bar.fill.solid()
        summary_bar.fill.fore_color.rgb = RGBColor(220, 20, 60)
        
        return slide
    
    def _extract_cost_data_from_llm(self, products: List[Dict[str, Any]], slide_content: Dict[str, Any]) -> Dict[str, Any]:
        """LLMデータからコスト情報を抽出して構造化"""
        try:
            # 商品情報から価格を抽出
            total_hardware_cost = 0
            for product in products:
                price_text = product.get('price', '0')
                # 価格から数値を抽出
                price_match = re.search(r'¥([0-9,]+)', price_text)
                if price_match:
                    price_str = price_match.group(1).replace(',', '')
                    try:
                        price = float(price_str)
                        total_hardware_cost += price
                    except ValueError:
                        pass
            
            # デフォルト値の設定
            current_data = {
                'equipment': '既存設備×100台',
                'hardware': '¥0',
                'running': '¥500,000/月',
                'electricity': '¥200,000/月',
                'tec': '150.00Kwh',
                'total': '¥700,000/月'
            }
            
            proposal_data = {
                'equipment': f'MPC×{len(products)}台',
                'hardware': 'カウンター料金に含む',
                'running': '¥300,000/月',
                'electricity': '¥120,000/月',
                'tec': '90.00Kwh',
                'total': '¥420,000/月'
            }
            
            # 差額計算
            current_monthly = 700000
            proposal_monthly = 420000
            monthly_diff = current_monthly - proposal_monthly
            
            difference_data = {
                'monthly': f'{monthly_diff:,}円',
                'yearly': f'{monthly_diff * 12:,}円',
                'five_year': f'{monthly_diff * 12 * 5:,}円'
            }
            
            # LLMデータから追加情報を抽出
            if slide_content:
                # 現状の課題から情報を抽出
                current_content = slide_content.get("slide1", {}).get("content", [])
                for line in current_content:
                    if "台" in line:
                        current_data['equipment'] = line.strip()
                    elif "円" in line and "月" in line:
                        current_data['running'] = line.strip()
                    elif "Kwh" in line:
                        current_data['tec'] = line.strip()
                
                # 導入メリットから情報を抽出
                benefits_content = slide_content.get("slide3", {}).get("content", [])
                for line in benefits_content:
                    if "削減" in line and "円" in line:
                        # 削減率から計算
                        reduction_match = re.search(r'([0-9]+)%', line)
                        if reduction_match:
                            reduction_rate = int(reduction_match.group(1)) / 100
                            proposal_data['running'] = f'¥{int(current_monthly * (1 - reduction_rate)):,}/月'
                            proposal_data['electricity'] = f'¥{int(200000 * (1 - reduction_rate)):,}/月'
                            break
                
                # トータルコストから情報を抽出
                total_content = slide_content.get("slide4", {}).get("content", [])
                for line in total_content:
                    if "総投資額" in line:
                        investment_match = re.search(r'¥([0-9,]+)', line)
                        if investment_match:
                            proposal_data['hardware'] = f'¥{investment_match.group(1)}'
                    elif "年間" in line and "削減" in line:
                        yearly_match = re.search(r'([0-9]+)%', line)
                        if yearly_match:
                            yearly_rate = int(yearly_match.group(1)) / 100
                            proposal_data['running'] = f'¥{int(current_monthly * (1 - yearly_rate)):,}/月'
            
            # 合計を再計算
            try:
                current_running = int(re.search(r'¥([0-9,]+)', current_data['running']).group(1).replace(',', ''))
                current_electricity = int(re.search(r'¥([0-9,]+)', current_data['electricity']).group(1).replace(',', ''))
                current_data['total'] = f'¥{current_running + current_electricity:,}/月'
                
                proposal_running = int(re.search(r'¥([0-9,]+)', proposal_data['running']).group(1).replace(',', ''))
                proposal_electricity = int(re.search(r'¥([0-9,]+)', proposal_data['electricity']).group(1).replace(',', ''))
                proposal_data['total'] = f'¥{proposal_running + proposal_electricity:,}/月'
                
                # 差額を再計算
                current_monthly = current_running + current_electricity
                proposal_monthly = proposal_running + proposal_electricity
                monthly_diff = current_monthly - proposal_monthly
                
                difference_data = {
                    'monthly': f'{monthly_diff:,}円',
                    'yearly': f'{monthly_diff * 12:,}円',
                    'five_year': f'{monthly_diff * 12 * 5:,}円'
                }
            except:
                pass
            
            return {
                'current': current_data,
                'proposal': proposal_data,
                'difference': difference_data
            }
            
        except Exception as e:
            # エラーが発生した場合、デフォルト値を返す
            return {
                'current': {
                    'equipment': '既存設備×100台',
                    'hardware': '¥0',
                    'running': '¥500,000/月',
                    'electricity': '¥200,000/月',
                    'tec': '150.00Kwh',
                    'total': '¥700,000/月'
                },
                'proposal': {
                    'equipment': f'MPC×{len(products)}台',
                    'hardware': 'カウンター料金に含む',
                    'running': '¥300,000/月',
                    'electricity': '¥120,000/月',
                    'tec': '90.00Kwh',
                    'total': '¥420,000/月'
                },
                'difference': {
                    'monthly': '280,000円',
                    'yearly': '3,360,000円',
                    'five_year': '16,800,000円'
                }
            }
    
    def _create_products_slide(self, prs: Presentation, title: str, content: List[str], products: List[Dict[str, Any]]) -> Any:
        """商品とその画像を含むスライドを作成"""
        # 空白レイアウトを使用
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # 会社ロゴを追加
        self._add_logo_to_slide(slide)
        
        # スライドタイトル
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        # 各商品のコンテンツを作成
        y_position = 1.8
        for i, product in enumerate(products):
            # 商品名
            product_title = slide.shapes.add_textbox(Inches(0.5), Inches(y_position), Inches(8), Inches(0.8))
            title_frame = product_title.text_frame
            title_frame.text = f"{i+1}. {product['name']} ({product['category']})"
            title_frame.paragraphs[0].font.size = Pt(20)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # 価格
            price_text = slide.shapes.add_textbox(Inches(0.5), Inches(y_position + 0.8), Inches(8), Inches(0.5))
            price_frame = price_text.text_frame
            price_frame.text = f"価格: {product['price']}"
            price_frame.paragraphs[0].font.size = Pt(16)
            price_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # 説明
            if product.get('description'):
                desc_text = slide.shapes.add_textbox(Inches(0.5), Inches(y_position + 1.3), Inches(8), Inches(1.2))
                desc_frame = desc_text.text_frame
                desc_frame.text = product['description'][:200] + ("..." if len(product['description']) > 200 else "")
                desc_frame.paragraphs[0].font.size = Pt(14)
                desc_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # 商品画像（右側）
            try:
                # プロジェクトルートからの相対パスを使用
                project_root = Path(__file__).parent.parent.parent
                image_path = project_root / "data/images/example_picture.png"
                if image_path.exists():
                    # テキストの右側に画像を配置
                    img_left = Inches(9.5)
                    img_top = Inches(y_position)
                    img_width = Inches(3.0)
                    img_height = Inches(2.0)
                    
                    slide.shapes.add_picture(str(image_path), img_left, img_top, img_width, img_height)
            except Exception as e:
                # 画像の追加に失敗した場合、プレースホルダーを追加
                placeholder = slide.shapes.add_textbox(img_left, img_top, img_width, img_height)
                placeholder_frame = placeholder.text_frame
                placeholder_frame.text = "画像\nなし"
                placeholder_frame.paragraphs[0].font.size = Pt(16)
                placeholder_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
                placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 次の商品の位置を増加
            y_position += 2.8
        
        return slide
    
    def _add_logo_to_slide(self, slide: Any):
        """スライドに会社ロゴを追加"""
        try:
            # プロジェクトルートからの相対パスを使用
            project_root = Path(__file__).parent.parent.parent
            logo_path = project_root / "data/images/otsuka_logo.jpg"
            if logo_path.exists():
                # 右上角にロゴを追加
                # ロゴの比例を保持 - サイズのみ縮小
                left = Inches(10.5)
                top = Inches(0.2)
                width = Inches(1.5)  # サイズを縮小するが、比例は保持
                
                # 元の画像サイズを取得して比例を保持
                from PIL import Image
                with Image.open(logo_path) as img:
                    orig_width, orig_height = img.size
                    aspect_ratio = orig_height / orig_width
                    height = width * aspect_ratio
                
                slide.shapes.add_picture(str(logo_path), left, top, width, height)
        except Exception as e:
            st.warning(f"ロゴの追加に失敗しました: {str(e)}")


# グローバルインスタンスを作成
slide_generator = SlideGenerator()
